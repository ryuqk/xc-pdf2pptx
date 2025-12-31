
import argparse
import os
import io
import fitz  # pymupdf
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from google import genai
from google.genai import types
import json
import base64
from dotenv import load_dotenv

load_dotenv()


from pptx.dml.color import RGBColor

class DocumentProcessor:
    def __init__(self, file_path):
        self.doc = fitz.open(file_path)

    def get_page_image(self, page_num, zoom=2.0):
        """Renders a PDF page to a PIL Image."""
        page = self.doc.load_page(page_num)
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        img_data = pix.tobytes("png")
        return Image.open(io.BytesIO(img_data)), page.rect.width, page.rect.height

    def close(self):
        self.doc.close()

class GeminiAnalyzer:
    def __init__(self, api_key):
        self.client = genai.Client(api_key=api_key)

    def analyze_page(self, image):
        """Analyzes the image using Gemini to identify text and figures."""
        
        # Convert PIL Image to bytes
        img_byte_arr = io.BytesIO()
        image.save(img_byte_arr, format='PNG')
        img_byte_arr = img_byte_arr.getvalue()

        prompt = """
        Analyze this document page image. I want to convert this into an editable PowerPoint slide.
        Identify two types of elements:
        1. "text_blocks": Select all visible text. Group related text (like paragraphs) together. 
           CRITICAL: The bounding box must be TIGHT around the text content. Do not include excessive empty space.
           For each block, provide:
           - "text": The actual text content.
           - "box_2d": The bounding box [ymin, xmin, ymax, xmax] normalized to 1000x1000.
           - "font_size_pt": Estimate the font size in points (approximate).
           - "font_color_hex": Estimate the font color in hex format (e.g. #000000).
           - "font_family": Enum "sans" (like Arial, Gothic) or "serif" (like Times, Mincho).
           - "is_bold": Boolean, true if the text is bold.
           - "is_title": Boolean, true if it looks like a title/heading.
        
        2. "image_regions": Identify non-text graphical elements (figures, diagrams, photos, icons, complex background shapes that generally shouldn't be executed as editable text). Do NOT include simple background colors or simple separators if possible, but do include main visual content.
           - "box_2d": The bounding box [ymin, xmin, ymax, xmax] normalized to 1000x1000.
           - "description": Short description of the image.
           
        Output strictly JSON format.
        """

        response = self.client.models.generate_content(
            model='gemini-3-flash-preview',
            contents=[
                types.Content(
                    role="user",
                    parts=[
                        types.Part.from_text(text=prompt),
                        types.Part.from_bytes(data=img_byte_arr, mime_type="image/png")
                    ]
                )
            ],
            config=types.GenerateContentConfig(
                response_mime_type="application/json",
                response_schema=types.Schema(
                    type=types.Type.OBJECT,
                    properties={
                        "text_blocks": types.Schema(
                            type=types.Type.ARRAY,
                            items=types.Schema(
                                type=types.Type.OBJECT,
                                properties={
                                    "text": types.Schema(type=types.Type.STRING),
                                    "box_2d": types.Schema(
                                        type=types.Type.ARRAY,
                                        items=types.Schema(type=types.Type.INTEGER)
                                    ),
                                    "font_size_pt": types.Schema(type=types.Type.NUMBER),
                                    "font_color_hex": types.Schema(type=types.Type.STRING),
                                    "font_family": types.Schema(type=types.Type.STRING),
                                    "is_bold": types.Schema(type=types.Type.BOOLEAN),
                                    "is_title": types.Schema(type=types.Type.BOOLEAN),
                                },
                            ),
                        ),
                        "image_regions": types.Schema(
                            type=types.Type.ARRAY,
                            items=types.Schema(
                                type=types.Type.OBJECT,
                                properties={
                                    "box_2d": types.Schema(
                                        type=types.Type.ARRAY,
                                        items=types.Schema(type=types.Type.INTEGER)
                                    ),
                                    "description": types.Schema(type=types.Type.STRING),
                                },
                            ),
                        ),
                    },
                ),
            )
        )
        
        try:
             # Depending on SDK version, response.text might be the JSON string
            return json.loads(response.text)
        except Exception as e:
            print(f"Error parsing Gemini response: {e}")
            print(f"Raw response: {response.text}")
            return {"text_blocks": [], "image_regions": []}

class PPTXBuilder:
    def __init__(self, output_path, mode="standard", font_scale=1.1):
        self.prs = Presentation()
        self.output_path = output_path
        self.mode = mode
        self.font_scale = font_scale
        
    def set_slide_size(self, width, height):
        self.prs.slide_width = int(width * 72 * 12700) 
        self.prs.slide_height = int(height * 72 * 12700)

    def get_edge_color(self, image, box_2d):
        """
        Calculates background color by sampling the edges of the region.
        box_2d: [ymin, xmin, ymax, xmax] normalized 1000
        """
        ymin, xmin, ymax, xmax = box_2d
        w, h = image.size
        
        # Expand slightly to capture background context if possible (ensure we are outside the letter forms)
        # But be careful not to go out of bounds or into other text.
        # Actually, using the box itself and checking edges is usually safe as text doesn't touch the box edge typically.
        
        left = int((xmin / 1000.0) * w)
        top = int((ymin / 1000.0) * h)
        right = int((xmax / 1000.0) * w)
        bottom = int((ymax / 1000.0) * h)
        
        # Clamp
        left = max(0, left)
        top = max(0, top)
        right = min(w, right)
        bottom = min(h, bottom)
        
        if right <= left or bottom <= top:
            return (255, 255, 255) 

        cropped = image.crop((left, top, right, bottom))
        
        # Sample edges
        pixels = []
        cw, ch = cropped.size
        
        # Sample logic: Top and Bottom rows, Left and Right columns
        # Take 2px depth if possible
        depth = 2
        
        try:
             # Convert to RGB to ensure tuples
            cropped_rgb = cropped.convert("RGB")
            
            # Top & Bottom
            for y in range(min(depth, ch)):
                for x in range(cw):
                    pixels.append(cropped_rgb.getpixel((x, y)))
                    pixels.append(cropped_rgb.getpixel((x, ch - 1 - y)))
            
            # Left & Right
            for x in range(min(depth, cw)):
                for y in range(ch):
                    pixels.append(cropped_rgb.getpixel((x, y)))
                    pixels.append(cropped_rgb.getpixel((cw - 1 - x, y)))
            
            if not pixels:
                return (255, 255, 255)
                
            # Find most common color
            from collections import Counter
            most_common = Counter(pixels).most_common(1)[0][0]
            return most_common
            
        except Exception:
            return (255, 255, 255)

    def add_slide(self, original_image, layout_data, pdf_width, pdf_height):
        # Create a blank slide
        blank_slide_layout = self.prs.slide_layouts[6] 
        slide = self.prs.slides.add_slide(blank_slide_layout)

        # Gemini uses 1000x1000 normalization
        scale_x = self.prs.slide_width / 1000.0
        scale_y = self.prs.slide_height / 1000.0
        
        if self.mode == "text_focus":
            # 1. Set background image (Full Page)
            img_stream = io.BytesIO()
            original_image.save(img_stream, format="PNG")
            img_stream.seek(0)
            
            slide.shapes.add_picture(img_stream, 0, 0, self.prs.slide_width, self.prs.slide_height)

        # 2. Add Images (Standard Mode only)
        if self.mode == "standard" and "image_regions" in layout_data:
            for img_region in layout_data["image_regions"]:
                ymin, xmin, ymax, xmax = img_region["box_2d"]
                
                w, h = original_image.size
                
                left = int((xmin / 1000.0) * w)
                top = int((ymin / 1000.0) * h)
                right = int((xmax / 1000.0) * w)
                bottom = int((ymax / 1000.0) * h)
                
                if right > left and bottom > top:
                    cropped_img = original_image.crop((left, top, right, bottom))
                    img_stream = io.BytesIO()
                    cropped_img.save(img_stream, format="PNG")
                    img_stream.seek(0)
                    
                    slide_left = int(xmin * scale_x)
                    slide_top = int(ymin * scale_y)
                    slide_width = int((xmax - xmin) * scale_x)
                    slide_height = int((ymax - ymin) * scale_y)
                    
                    try:
                        slide.shapes.add_picture(img_stream, slide_left, slide_top, slide_width, slide_height)
                    except Exception as e:
                        print(f"Failed to add image: {e}")

        # 3. Add Text
        if "text_blocks" in layout_data:
            for text_block in layout_data["text_blocks"]:
                ymin_norm, xmin_norm, ymax_norm, xmax_norm = text_block["box_2d"]
                text_content = text_block.get("text", "")
                font_size = text_block.get("font_size_pt", 12)
                font_color_hex = text_block.get("font_color_hex", "#000000")
                font_family_style = text_block.get("font_family", "sans")
                is_bold = text_block.get("is_bold", False)
                
                # Coordinates
                left = int(xmin_norm * scale_x)
                top = int(ymin_norm * scale_y)
                width = int((xmax_norm - xmin_norm) * scale_x)
                height = int((ymax_norm - ymin_norm) * scale_y)

                if self.mode == "text_focus":
                    # --- Strategy: Two Shapes ---
                    # Shape 1: Mask Rectangle (Inflated, Filled with Bg Color, No Text)
                    # Shape 2: Text Box (Standard Coordinates, No Fill, Text)
                    
                    # 1. MASK SHAPE
                    # Sample color from the original box (tight) edges to get accurate background
                    bg_color = self.get_edge_color(original_image, [ymin_norm, xmin_norm, ymax_norm, xmax_norm])
                    
                    # Inflate box for masking
                    # 5 units out of 1000 approx 0.5%
                    inflation = 5 
                    mask_ymin = max(0, ymin_norm - inflation)
                    mask_xmin = max(0, xmin_norm - inflation)
                    mask_ymax = min(1000, ymax_norm + inflation)
                    mask_xmax = min(1000, xmax_norm + inflation)

                    mask_left = int(mask_xmin * scale_x)
                    mask_top = int(mask_ymin * scale_y)
                    mask_width = int((mask_xmax - mask_xmin) * scale_x)
                    mask_height = int((mask_ymax - mask_ymin) * scale_y)
                    
                    mask_shape = slide.shapes.add_shape(
                        1, # MSO_SHAPE.RECTANGLE (1 is value)
                        mask_left, mask_top, mask_width, mask_height
                    )
                    mask_shape.fill.solid()
                    mask_shape.fill.fore_color.rgb = RGBColor(bg_color[0], bg_color[1], bg_color[2])
                    mask_shape.line.fill.background() # No line
                    
                    # Remove shadow
                    mask_shape.shadow.inherit = False 
                    
                    # 2. TEXT SHAPE (Precise)
                    # Use original coordinates
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    
                    # Set margins to 0 for tight alignment
                    tf = txBox.text_frame
                    tf.margin_left = 0
                    tf.margin_right = 0
                    tf.margin_top = 0
                    tf.margin_bottom = 0
                    tf.word_wrap = True 

                else:
                    # Standard Mode
                    txBox = slide.shapes.add_textbox(left, top, width, height)
                    tf = txBox.text_frame
                    tf.word_wrap = True 

                tf.text = text_content
                
                # Apply styles to ALL paragraphs to handle multi-line text correctly
                for p in tf.paragraphs:
                    # Font Size (Scaled)
                    scaled_font_size = (font_size * self.font_scale) if font_size else 12
                    p.font.size = Pt(scaled_font_size)
                    
                    # Bold
                    if is_bold:
                        p.font.bold = True
                    
                    # Font Family logic
                    if self.mode == "text_focus":
                        if font_family_style == "serif":
                             p.font.name = "MS Mincho"
                        else:
                             p.font.name = "Meiryo" 

                    # Color logic
                    try:
                        if font_color_hex and font_color_hex.startswith("#"):
                            r = int(font_color_hex[1:3], 16)
                            g = int(font_color_hex[3:5], 16)
                            b = int(font_color_hex[5:7], 16)
                            p.font.color.rgb = RGBColor(r, g, b)
                    except:
                        pass

    def save(self):
        self.prs.save(self.output_path)
        print(f"Presentation saved to {self.output_path}")

def main():
    parser = argparse.ArgumentParser(description="Convert PDF or Images to editable PPTX using Gemini.")
    parser.add_argument("input_file", help="Path to input file (PDF, PNG, JPG, etc.)")
    parser.add_argument("output_pptx", help="Path to output PPTX file")
    parser.add_argument("--api_key", help="Google Gemini API Key", default=os.environ.get("GOOGLE_API_KEY"))
    parser.add_argument("--mode", help="Conversion mode: 'standard' or 'text_focus'", default="standard", choices=["standard", "text_focus"])
    parser.add_argument("--font_scale", help="Font size scaling factor", default=1.1, type=float)
    
    args = parser.parse_args()
    
    if not args.api_key:
        print("Error: API Key is required. Set GOOGLE_API_KEY env var or pass --api_key.")
        return

    if not os.path.exists(args.input_file):
        print(f"Error: Input file {args.input_file} not found.")
        return

    proc = DocumentProcessor(args.input_file)
    analyzer = GeminiAnalyzer(args.api_key)
    builder = PPTXBuilder(args.output_pptx, mode=args.mode, font_scale=args.font_scale)
    
    print(f"Processing {args.input_file} in {args.mode} mode...")

    try:
        for page_num in range(len(proc.doc)):
            print(f"Processing page {page_num + 1}/{len(proc.doc)}...")
            image, w, h = proc.get_page_image(page_num)
            
            if page_num == 0:
                 builder.prs.slide_width = int(w * 12700)
                 builder.prs.slide_height = int(h * 12700)

            layout_data = analyzer.analyze_page(image)
            builder.add_slide(image, layout_data, w, h)
            
    finally:
        proc.close()
        
    builder.save()


if __name__ == "__main__":
    main()
